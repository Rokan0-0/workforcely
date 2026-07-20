from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions (Workforcely Theme)
DARK_BG = RGBColor(15, 23, 42)       # Slate 900 #0F172A
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50 #F8FAFC
PRIMARY = RGBColor(30, 58, 95)        # Deep Navy #1E3A5F
ACCENT = RGBColor(20, 184, 166)       # Teal Accent #14B8A6
SUCCESS = RGBColor(16, 185, 129)     # Emerald #10B981
CARD_BG = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_LIGHT = RGBColor(255, 255, 255)
TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
BORDER_COLOR = RGBColor(226, 232, 240)

def set_slide_background(slide, color):
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = color
    bg_rect.line.fill.background()

def add_header(slide, title_text, category_text=None, is_dark=False):
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Segoe UI"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT if is_dark else TEXT_DARK

def add_footer(slide, current_page, total_pages, is_dark=False):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Workforcely HR | Training LMS & HMO Healthcare Portal"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    
    p_num = tf.add_paragraph()
    p_num.text = f"Slide {current_page} of {total_pages}"
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = "Segoe UI"
    p_num.font.size = Pt(9)
    p_num.font.color.rgb = TEXT_MUTED

# ==========================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==========================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s1, DARK_BG)

t_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = t_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "Workforcely"
p1.font.name = "Segoe UI"
p1.font.size = Pt(60)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT

p2 = tf.add_paragraph()
p2.text = "AI Learning Management System & HMO Healthcare Benefits"
p2.font.name = "Segoe UI"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = ACCENT

p3 = tf.add_paragraph()
p3.text = "\nAutomating Employee Upskilling & Healthcare Self-Service for Modern Teams"
p3.font.name = "Segoe UI"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

add_footer(s1, 1, 5, is_dark=True)

# ==========================================================
# SLIDE 2: AI-Powered LMS (Light Theme)
# ==========================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s2, LIGHT_BG)
add_header(s2, "AI-Powered Learning Management System (LMS)", "Upskilling & Training")

# Left Column (Features)
left_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
tf2 = left_box.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "Enterprise Learning & Training Studio:"
p.font.name = "Segoe UI"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_DARK

features_lms = [
    ("🤖 AI Course Co-pilot", "HR types a prompt (e.g., 'Data Privacy & NDPR Compliance') and Gemini AI generates the full course title, markdown lessons, and quiz questions in seconds."),
    ("🎨 Course Creation Studio", "4-Step Wizard: Basic Details → Curriculum Editor (Markdown & Video Embeds) → Quiz Authoring (Pass Marks & Retakes) → Live Student Preview."),
    ("📚 Employee Learning Portal", "Self-service course catalog filterable by category and difficulty. Interactive lesson viewer with dark/light mode readability."),
    ("🏆 Automated PDF/HTML Certificates", "Instant generation of signed Certificates of Completion upon passing the course evaluation quiz.")
]

for title, desc in features_lms:
    p_t = tf2.add_paragraph()
    p_t.text = f"\n{title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = PRIMARY

    p_d = tf2.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

# Right Card (Summary Stats)
right_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8))
right_card.fill.solid()
right_card.fill.fore_color.rgb = CARD_BG
right_card.line.color.rgb = BORDER_COLOR

rc_tf = right_card.text_frame
rc_tf.word_wrap = True

rp1 = rc_tf.paragraphs[0]
rp1.text = "LMS IMPACT HIGHLIGHTS"
rp1.font.name = "Segoe UI"
rp1.font.size = Pt(11)
rp1.font.bold = True
rp1.font.color.rgb = ACCENT

rp2 = rc_tf.add_paragraph()
rp2.text = "\n1 Click"
rp2.font.name = "Segoe UI"
rp2.font.size = Pt(48)
rp2.font.bold = True
rp2.font.color.rgb = PRIMARY

rp3 = rc_tf.add_paragraph()
rp3.text = "AI course generation time for HR admins."
rp3.font.name = "Segoe UI"
rp3.font.size = Pt(12)
rp3.font.color.rgb = TEXT_MUTED

rp4 = rc_tf.add_paragraph()
rp4.text = "\n100% Automated"
rp4.font.name = "Segoe UI"
rp4.font.size = Pt(36)
rp4.font.bold = True
rp4.font.color.rgb = SUCCESS

rp5 = rc_tf.add_paragraph()
rp5.text = "Quiz grading and instant certificate issuance."
rp5.font.name = "Segoe UI"
rp5.font.size = Pt(12)
rp5.font.color.rgb = TEXT_MUTED

add_footer(s2, 2, 5)

# ==========================================================
# SLIDE 3: HMO Healthcare Portal (Light Theme)
# ==========================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s3, LIGHT_BG)
add_header(s3, "HMO & Healthcare Benefits Portal", "Employee Well-being")

# Left Column (Features)
left_box3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
tf3 = left_box3.text_frame
tf3.word_wrap = True

p = tf3.paragraphs[0]
p.text = "Automated Health Insurance Management:"
p.font.name = "Segoe UI"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEXT_DARK

features_hmo = [
    ("🏥 HMO Provider Tiers", "Support for leading Nigerian providers (Reliance HMO, Hygeia, AXA Mansard, Leadway) across Bronze, Silver, Gold, and Platinum tiers."),
    ("👨‍👩‍👧‍👦 Family Dependant Self-Service", "Employees register spouse and children directly on their profile, eliminating paper forms and HR email back-and-forth."),
    ("🏨 Nationwide Hospital Access", "Displays total covered hospitals (up to 1,800+ nationwide) and annual financial coverage limits per plan."),
    ("📊 HR Cost & Enrollee Oversight", "Real-time HR dashboard tracking monthly healthcare costs, active enrollee counts, and dependant metrics.")
]

for title, desc in features_hmo:
    p_t = tf3.add_paragraph()
    p_t.text = f"\n{title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT

    p_d = tf3.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

# Right Card (HMO Tiers Showcase)
hmo_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8))
hmo_card.fill.solid()
hmo_card.fill.fore_color.rgb = CARD_BG
hmo_card.line.color.rgb = BORDER_COLOR

htf = hmo_card.text_frame
htf.word_wrap = True

hp1 = htf.paragraphs[0]
hp1.text = "SUPPORTED HMO TIERS"
hp1.font.name = "Segoe UI"
hp1.font.size = Pt(11)
hp1.font.bold = True
hp1.font.color.rgb = PRIMARY

tiers = [
    ("Bronze Tier", "Reliance HMO • 450+ Hospitals • ₦1.2M Limit"),
    ("Silver Tier", "Hygeia HMO • 850+ Hospitals • ₦2.5M Limit"),
    ("Gold Tier", "AXA Mansard • 1,200+ Hospitals • ₦5.0M Limit"),
    ("Platinum Tier", "Leadway Health • 1,800+ Hospitals • ₦10.0M Limit")
]

for tier_name, tier_info in tiers:
    tp1 = htf.add_paragraph()
    tp1.text = f"\n• {tier_name}"
    tp1.font.name = "Segoe UI"
    tp1.font.size = Pt(14)
    tp1.font.bold = True
    tp1.font.color.rgb = ACCENT
    
    tp2 = htf.add_paragraph()
    tp2.text = f"  {tier_info}"
    tp2.font.name = "Segoe UI"
    tp2.font.size = Pt(11)
    tp2.font.color.rgb = TEXT_MUTED

add_footer(s3, 3, 5)

# ==========================================================
# SLIDE 4: Business Value & Synergy (Light Theme)
# ==========================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s4, LIGHT_BG)
add_header(s4, "The Value Proposition: Why LMS + HMO Matter", "Business Impact")

card_w = Inches(5.6)
card_h = Inches(4.5)
gap = Inches(0.5)

# Card 1: LMS Impact
c1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), card_w, card_h)
c1.fill.solid()
c1.fill.fore_color.rgb = CARD_BG
c1.line.color.rgb = PRIMARY
c1.line.width = Pt(2)

c1_tf = c1.text_frame
c1_tf.word_wrap = True

cp1 = c1_tf.paragraphs[0]
cp1.text = "LMS: SKILL & COMPLIANCE IMPACT"
cp1.font.name = "Segoe UI"
cp1.font.size = Pt(13)
cp1.font.bold = True
cp1.font.color.rgb = PRIMARY

lms_impacts = [
    "Zero Content Bottleneck: AI Co-pilot creates courses instantly.",
    "Audit Readiness: Tracks mandatory compliance training completion.",
    "Employee Engagement: Clear learning paths & verifiable certificates.",
    "Scalable Onboarding: New hires complete training on day 1."
]
for imp in lms_impacts:
    p = c1_tf.add_paragraph()
    p.text = f"\n✓ {imp}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK

# Card 2: HMO Impact
c2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), card_w, card_h)
c2.fill.solid()
c2.fill.fore_color.rgb = CARD_BG
c2.line.color.rgb = ACCENT
c2.line.width = Pt(2)

c2_tf = c2.text_frame
c2_tf.word_wrap = True

cp2 = c2_tf.paragraphs[0]
cp2.text = "HMO: RETENTION & WELL-BEING"
cp2.font.name = "Segoe UI"
cp2.font.size = Pt(13)
cp2.font.bold = True
cp2.font.color.rgb = ACCENT

hmo_impacts = [
    "Higher Staff Retention: Comprehensive family healthcare benefits.",
    "Zero Paperwork: Direct online dependant registration.",
    "Cost Transparency: Real-time HR monthly budget overview.",
    "Peace of Mind: Instant access to nationwide hospital networks."
]
for imp in hmo_impacts:
    p = c2_tf.add_paragraph()
    p.text = f"\n✓ {imp}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK

add_footer(s4, 4, 5)

# ==========================================================
# SLIDE 5: Conclusion & Call To Action (Dark Theme)
# ==========================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s5, DARK_BG)

t_box5 = s5.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf5 = t_box5.text_frame
tf5.word_wrap = True

p1 = tf5.paragraphs[0]
p1.text = "Empower Your Workforce with Workforcely"
p1.font.name = "Segoe UI"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT

p2 = tf5.add_paragraph()
p2.text = "AI-Powered Learning + Modern Healthcare Benefits for African SMEs"
p2.font.name = "Segoe UI"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT

p3 = tf5.add_paragraph()
p3.text = "\nExperience the Live Platform:\n• Training Studio: /dashboard/training\n• Healthcare Portal: /dashboard/benefits"
p3.font.name = "Segoe UI"
p3.font.size = Pt(13)
p3.font.color.rgb = TEXT_MUTED

add_footer(s5, 5, 5, is_dark=True)

# Save
output_path = "Workforcely_LMS_and_HMO_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
