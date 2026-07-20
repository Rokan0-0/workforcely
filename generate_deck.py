from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
DARK_BG = RGBColor(15, 23, 42)       # Slate 900
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
PRIMARY = RGBColor(99, 102, 241)     # Indigo 500
PRIMARY_DARK = RGBColor(79, 70, 229) # Indigo 600
SECONDARY = RGBColor(59, 130, 246)   # Blue 500
SUCCESS = RGBColor(16, 185, 129)     # Emerald 500
WARNING = RGBColor(245, 158, 11)     # Amber 500
DANGER = RGBColor(239, 68, 68)       # Red 500
CARD_BG = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_LIGHT = RGBColor(255, 255, 255)
TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
BORDER_COLOR = RGBColor(226, 232, 240)# Slate 200

def set_slide_background(slide, color):
    """Draws a full-screen rectangle to serve as a reliable background color."""
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = color
    bg_rect.line.fill.background()
    # Send to back is handled implicitly as it's the first shape added

def add_header(slide, title_text, category_text=None, is_dark=False):
    """Adds a standard premium header to the slide."""
    # Category Tracker
    if category_text:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Segoe UI"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY if not is_dark else PRIMARY

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT if is_dark else TEXT_DARK

def add_footer(slide, current_page, total_pages, is_dark=False):
    """Adds page number and branding to the bottom of the slide."""
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Workforcely | Nigerian SME HR & LMS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED if not is_dark else TEXT_MUTED
    
    # Page Number
    p_num = tf.add_paragraph()
    p_num.text = f"Slide {current_page} of {total_pages}"
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = "Segoe UI"
    p_num.font.size = Pt(9)
    p_num.font.color.rgb = TEXT_MUTED if not is_dark else TEXT_MUTED

def add_placeholder_or_image(slide, x, y, width, height, image_name, desc_text):
    """Inserts a screenshot if present, otherwise draws a premium colored placeholder card."""
    if os.path.exists(image_name):
        slide.shapes.add_picture(image_name, x, y, width, height)
    else:
        # Draw placeholder card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(241, 245, 249) # Light slate
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Inner text frame
        tb = slide.shapes.add_textbox(x, y + (height / 2) - Inches(0.7), width, Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"[SCREENSHOT PLACEHOLDER]"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        
        p2 = tf.add_paragraph()
        p2.text = f"Save screenshot as: {image_name}\n({desc_text})"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED

# ==========================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==========================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s1, DARK_BG)

# Title Text Frame
t_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = t_box.text_frame
tf.word_wrap = True

# Main Title
p1 = tf.paragraphs[0]
p1.text = "Workforcely"
p1.font.name = "Segoe UI"
p1.font.size = Pt(64)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT

# Subtitle
p2 = tf.add_paragraph()
p2.text = "The Operational Engine for African SMEs"
p2.font.name = "Segoe UI"
p2.font.size = Pt(24)
p2.font.color.rgb = PRIMARY

# Details
p3 = tf.add_paragraph()
p3.text = "\nConsolidating the Employee Lifecycle, Localized Compliance & AI Operations"
p3.font.name = "Segoe UI"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

add_footer(s1, 1, 10, is_dark=True)

# ==========================================================
# SLIDE 2: The Spreadsheet Nightmare (Light Theme / Problem)
# ==========================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s2, LIGHT_BG)
add_header(s2, "The Spreadsheet & Compliance Nightmare", "Problem Statement")

# Large Left Text Box (The Challenge)
chal_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
tf_chal = chal_box.text_frame
tf_chal.word_wrap = True

p = tf_chal.paragraphs[0]
p.text = "Managing staff is the #1 operational bottleneck for Nigerian SMEs:"
p.font.name = "Segoe UI"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = TEXT_DARK

bullets = [
    ("Manual Tax Calculations", "Calculating monthly progressive PAYE taxes and pension deductions (8% of basic, housing, transport) by hand leads to persistent compliance errors."),
    ("Fragmented Operations", "Using WhatsApp groups for leaves, paper files for folders, and Excel for attendance leads to high friction and administrative overhead."),
    ("Audit Fine Risks", "Failing audits due to poor record-keeping triggers steep government penalties that threaten small business survival.")
]

for title, desc in bullets:
    p_t = tf_chal.add_paragraph()
    p_t.text = f"\n• {title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = DANGER
    
    p_d = tf_chal.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED

# Right Side Card (The Statistics)
stat_card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
stat_card.fill.solid()
stat_card.fill.fore_color.rgb = CARD_BG
stat_card.line.color.rgb = BORDER_COLOR

stat_tf = stat_card.text_frame
stat_tf.word_wrap = True

sp1 = stat_tf.paragraphs[0]
sp1.text = "THE COST OF INEFFICIENCY"
sp1.font.name = "Segoe UI"
sp1.font.size = Pt(11)
sp1.font.bold = True
sp1.font.color.rgb = PRIMARY_DARK

sp2 = stat_tf.add_paragraph()
sp2.text = "\n15 Days"
sp2.font.name = "Segoe UI"
sp2.font.size = Pt(54)
sp2.font.bold = True
sp2.font.color.rgb = DANGER

sp3 = stat_tf.add_paragraph()
sp3.text = "Wasted annually by a typical SME owner on HR paperwork & manual calculations."
sp3.font.name = "Segoe UI"
sp3.font.size = Pt(14)
sp3.font.color.rgb = TEXT_DARK

sp4 = stat_tf.add_paragraph()
sp4.text = "\n39 Million"
sp4.font.name = "Segoe UI"
sp4.font.size = Pt(36)
sp4.font.bold = True
sp4.font.color.rgb = TEXT_DARK

sp5 = stat_tf.add_paragraph()
sp5.text = "Micro & small businesses in Nigeria facing these exact operational bottlenecks."
sp5.font.name = "Segoe UI"
sp5.font.size = Pt(12)
sp5.font.color.rgb = TEXT_MUTED

add_footer(s2, 2, 10)

# ==========================================================
# SLIDE 3: Meet Workforcely (Light Theme / Solution)
# ==========================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s3, LIGHT_BG)
add_header(s3, "Introducing Workforcely", "The Solution")

# Intro Subtitle
sub_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
sub_box.text_frame.word_wrap = True
sp = sub_box.text_frame.paragraphs[0]
sp.text = "An all-in-one operational engine designed specifically for growing retail, fintech, and service teams."
sp.font.name = "Segoe UI"
sp.font.size = Pt(16)
sp.font.color.rgb = TEXT_MUTED

# Three Pillars Layout (Cards)
card_width = Inches(3.64)
card_height = Inches(4.2)
gap = Inches(0.4)
y_offset = Inches(2.3)

pillars = [
    ("Employee Lifecycle", "Manages the entire staff journey. From candidate tracking, automatic profile instantiation, onboarding checklists, leaves, to performance calibration scorecards.", PRIMARY),
    ("Localized Compliance", "Calculates progressive Nigerian PAYE income tax, standard 8% pension deductions, and allows flexible HR overrides (bonuses, deductions) with mandatory reason-logging.", SECONDARY),
    ("AI-Powered Operations", "Integrates Google Gemini to query organization-wide data (leaves, payroll totals, attendance) or private self-service data via conversational queries.", SUCCESS)
]

for idx, (title, desc, color) in enumerate(pillars):
    x_offset = Inches(0.8) + idx * (card_width + gap)
    card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_offset, y_offset, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1)
    
    # Text
    tf = card.text_frame
    tf.word_wrap = True
    
    # Bullet Color Accent Line inside Card
    accent_bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_offset + Inches(0.3), y_offset + Inches(0.3), Inches(0.8), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = color
    accent_bar.line.fill.background()
    
    # Text Content
    tb = s3.shapes.add_textbox(x_offset + Inches(0.2), y_offset + Inches(0.6), card_width - Inches(0.4), card_height - Inches(0.8))
    tb_tf = tb.text_frame
    tb_tf.word_wrap = True
    
    p = tb_tf.paragraphs[0]
    p.text = f"\n{title}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    p_desc = tb_tf.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED

add_footer(s3, 3, 10)

# ==========================================================
# SLIDE 4: Self-Service & Local Taxes (Light Theme / LMS + Attendance)
# ==========================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s4, LIGHT_BG)
add_header(s4, "Self-Service, Attendance & Automated Tax", "Employee View")

# Left Column (Features)
left_box = s4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.5))
tf = left_box.text_frame
tf.word_wrap = True

p_intro = tf.paragraphs[0]
p_intro.text = "Empowering employees with self-service utility:"
p_intro.font.name = "Segoe UI"
p_intro.font.size = Pt(16)
p_intro.font.bold = True
p_intro.font.color.rgb = TEXT_DARK

features = [
    ("Mobile Punch Clock", "Employees clock in/out directly from their dashboard. Tracks check-in times and feeds automatically into performance scorecards."),
    ("Compliance Payslips", "Instant breakdowns of Basic, Housing, and Transport allowances alongside automatic Nigerian PAYE tax and 8% pension deductions."),
    ("Leave Requests", "Submits request reasons, logs dates, and displays real-time approval status with email notifications.")
]

for title, desc in features:
    p_t = tf.add_paragraph()
    p_t.text = f"\n• {title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = PRIMARY_DARK
    
    p_d = tf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

# Right Column (Image/Screenshot)
add_placeholder_or_image(
    s4, Inches(6.4), Inches(1.8), Inches(6.1), Inches(4.5),
    "screenshot_employee.png", "Fatima's employee dashboard or payslip page"
)

add_footer(s4, 4, 10)

# ==========================================================
# SLIDE 5: Hiring & Automating Onboarding (Light Theme / Kanban)
# ==========================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s5, LIGHT_BG)
add_header(s5, "Streamlining Recruitment & Onboarding", "HR Admin View")

# Left Column (Image/Screenshot)
add_placeholder_or_image(
    s5, Inches(0.8), Inches(1.8), Inches(6.1), Inches(4.5),
    "screenshot_recruitment.png", "Applicant Kanban board or onboarding checklists"
)

# Right Column (Features)
right_box = s5.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.5))
tf = right_box.text_frame
tf.word_wrap = True

p_intro = tf.paragraphs[0]
p_intro.text = "From candidate sourcing to day-one readiness:"
p_intro.font.name = "Segoe UI"
p_intro.font.size = Pt(16)
p_intro.font.bold = True
p_intro.font.color.rgb = TEXT_DARK

features = [
    ("Visual Kanban Board", "Drag candidates seamlessly through pipeline stages: Applied → Shortlisted → Interview → Offer → Hired/Rejected."),
    ("Auto-Profile Generation", "Moving a candidate to 'Hired' instantly instantiates their employee record, sets up login credentials, and drafts onboarding tasks."),
    ("Checklist Tracking", "Interactive checklists (signed contract, ID capture, account setup) updated by employee and audited by HR in real-time.")
]

for title, desc in features:
    p_t = tf.add_paragraph()
    p_t.text = f"\n• {title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = PRIMARY_DARK
    
    p_d = tf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

add_footer(s5, 5, 10)

# ==========================================================
# SLIDE 6: The AI HR Assistant (Light Theme / Gemini Assistant)
# ==========================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s6, LIGHT_BG)
add_header(s6, "AI HR Assistant: Instant Business Intelligence", "Secret Sauce")

# Left Column (Features)
left_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.5))
tf = left_box.text_frame
tf.word_wrap = True

p_intro = tf.paragraphs[0]
p_intro.text = "Conversational data access powered by Google Gemini:"
p_intro.font.name = "Segoe UI"
p_intro.font.size = Pt(16)
p_intro.font.bold = True
p_intro.font.color.rgb = TEXT_DARK

features = [
    ("Natural Language Queries", "Admins ask plain English questions: 'Who is on leave next week?' or 'Total payroll cost this month' instead of exporting CSVs."),
    ("Context-Aware Logic", "Gemini interprets request context, queries our relational tables, and returns clean, markdown-formatted tables or summaries."),
    ("Strict Data Scoping", "Self-service portal is sandboxed to the employee's personal files. Employees can only ask about their own leaves, checklists, or payslips.")
]

for title, desc in features:
    p_t = tf.add_paragraph()
    p_t.text = f"\n• {title}"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = SUCCESS
    
    p_d = tf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = "Segoe UI"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

# Right Column (Image/Screenshot)
add_placeholder_or_image(
    s6, Inches(6.4), Inches(1.8), Inches(6.1), Inches(4.5),
    "screenshot_ai.png", "AI assistant view with query: 'Who is on leave next week?'"
)

add_footer(s6, 6, 10)

# ==========================================================
# SLIDE 7: Market Opportunity (Light Theme)
# ==========================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s7, LIGHT_BG)
add_header(s7, "A Massive Target Market in West Africa", "Market Opportunity")

# Subtitle
sub_box = s7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
sp = sub_box.text_frame.paragraphs[0]
sp.text = "Capturing the operational software shift in emerging market small businesses."
sp.font.name = "Segoe UI"
sp.font.size = Pt(16)
sp.font.color.rgb = TEXT_MUTED

# Stat Grid (Three Big numbers)
stat_box_width = Inches(3.64)
stat_box_height = Inches(4.0)
gap = Inches(0.4)
y_offset = Inches(2.2)

stats = [
    ("39M+", "Nigerian SMEs", "Micro, small, and medium businesses represent 96% of active businesses and 84% of employment in Nigeria, yet 90% run on spreadsheets or paper."),
    ("$35B+", "Addressable Payroll", "Estimated annual salary transaction volume processed by formal and semi-formal SMEs that need localized tax and pension compliance."),
    ("10x", "Administration Speed", "Replacing fragmented setups (WhatsApp, local Excel files, paper calendars) with Workforcely drives operational efficiency.")
]

for idx, (stat, label, desc) in enumerate(stats):
    x_offset = Inches(0.8) + idx * (stat_box_width + gap)
    card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_offset, y_offset, stat_box_width, stat_box_height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    
    tf = card.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"\n{stat}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    
    p_lbl = tf.add_paragraph()
    p_lbl.text = f"{label}"
    p_lbl.font.name = "Segoe UI"
    p_lbl.font.size = Pt(16)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = TEXT_DARK
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED

add_footer(s7, 7, 10)

# ==========================================================
# SLIDE 8: Simple Pricing Tiers (Light Theme)
# ==========================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s8, LIGHT_BG)
add_header(s8, "Low-Friction Business Model", "Pricing Tiers")

# Intro Subtitle
sub_box = s8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
sp = sub_box.text_frame.paragraphs[0]
sp.text = "Virality-driven packaging designed to hook growing businesses early."
sp.font.name = "Segoe UI"
sp.font.size = Pt(16)
sp.font.color.rgb = TEXT_MUTED

# Two Core Pricing Columns
card_width = Inches(5.6)
card_height = Inches(4.2)
gap = Inches(0.5)
y_offset = Inches(2.2)

# Free Tier
card1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_offset, card_width, card_height)
card1.fill.solid()
card1.fill.fore_color.rgb = CARD_BG
card1.line.color.rgb = BORDER_COLOR
tf1 = card1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "\nFREE FOREVER"
p.font.name = "Segoe UI"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = TEXT_MUTED

p_pr = tf1.add_paragraph()
p_pr.text = "N0 / month"
p_pr.font.name = "Segoe UI"
p_pr.font.size = Pt(40)
p_pr.font.bold = True
p_pr.font.color.rgb = TEXT_DARK

p_sub = tf1.add_paragraph()
p_sub.text = "For teams up to 5 employees\n"
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(13)
p_sub.font.bold = True
p_sub.font.color.rgb = PRIMARY

features_free = [
    "Punch clock & attendance tracking",
    "Basic employee directory",
    "Self-service leave requests",
    "Local text-based course builder"
]
for f in features_free:
    p_f = tf1.add_paragraph()
    p_f.text = f"✓  {f}"
    p_f.font.name = "Segoe UI"
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = TEXT_MUTED

# Premium Tier
card2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + card_width + gap, y_offset, card_width, card_height)
card2.fill.solid()
card2.fill.fore_color.rgb = CARD_BG
card2.line.color.rgb = PRIMARY
card2.line.width = Pt(2.5) # Accent border for active tier
tf2 = card2.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "\nPREMIUM GROW"
p.font.name = "Segoe UI"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY_DARK

p_pr = tf2.add_paragraph()
p_pr.text = "N1,200"
p_pr.font.name = "Segoe UI"
p_pr.font.size = Pt(40)
p_pr.font.bold = True
p_pr.font.color.rgb = TEXT_DARK

p_sub = tf2.add_paragraph()
p_sub.text = "Per employee / month\n"
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(13)
p_sub.font.bold = True
p_sub.font.color.rgb = PRIMARY

features_prem = [
    "Nigerian PAYE & Pension computation engine",
    "Gemini AI HR Assistant Integration",
    "Recruitment pipelines & automated onboarding",
    "Goal weight scorecards & custom evaluations"
]
for f in features_prem:
    p_f = tf2.add_paragraph()
    p_f.text = f"✓  {f}"
    p_f.font.name = "Segoe UI"
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = TEXT_MUTED

add_footer(s8, 8, 10)

# ==========================================================
# SLIDE 9: Roadmap (Light Theme)
# ==========================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s9, LIGHT_BG)
add_header(s9, "Future Roadmap & Deep Integration", "Development Roadmap")

# 3 Roadmap Steps (Horizontal Flow)
step_width = Inches(3.64)
step_height = Inches(4.0)
gap = Inches(0.4)
y_offset = Inches(2.2)

roadmap_steps = [
    ("Phase 1: Foundations", "Current Milestone", "Delivered localized payroll computations, leave approval portals, training LMS, applicant Kanban boards, and a context-aware Gemini AI assistant.", SUCCESS),
    ("Phase 2: Banking APIs", "Next 6 Months", "Integrate with local Nigerian bank APIs (e.g. Mono, Providus, Flutterwave) to support direct, one-click bulk salary payouts from the dashboard.", PRIMARY),
    ("Phase 3: Automated Filing", "Next 12 Months", "Direct integration with LIRS (Lagos State Internal Revenue Service) and PenCom portals to support direct automated filing of monthly taxes and deductions.", SECONDARY)
]

for idx, (title, milestone, desc, color) in enumerate(roadmap_steps):
    x_offset = Inches(0.8) + idx * (step_width + gap)
    card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_offset, y_offset, step_width, step_height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    
    tf = card.text_frame
    tf.word_wrap = True
    
    # Bullet Line Indicator
    accent = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_offset + Inches(0.3), y_offset + Inches(0.3), Inches(0.8), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()
    
    tb = s9.shapes.add_textbox(x_offset + Inches(0.2), y_offset + Inches(0.5), step_width - Inches(0.4), step_height - Inches(0.7))
    tb_tf = tb.text_frame
    tb_tf.word_wrap = True
    
    p = tb_tf.paragraphs[0]
    p.text = f"\n{title}"
    p.font.name = "Segoe UI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    p_m = tb_tf.add_paragraph()
    p_m.text = f"{milestone}"
    p_m.font.name = "Segoe UI"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = color
    
    p_desc = tb_tf.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED

add_footer(s9, 9, 10)

# ==========================================================
# SLIDE 10: Call To Action (Dark Theme)
# ==========================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(s10, DARK_BG)

# Title Text Frame
t_box = s10.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = t_box.text_frame
tf.word_wrap = True

# Main Call to Action
p1 = tf.paragraphs[0]
p1.text = "Let's Build the Future of SME Operations"
p1.font.name = "Segoe UI"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT

# Subtext
p2 = tf.add_paragraph()
p2.text = "Empowering emerging market businesses to stay compliant, scale, and thrive."
p2.font.name = "Segoe UI"
p2.font.size = Pt(18)
p2.font.color.rgb = PRIMARY

# Contact / Call details
p3 = tf.add_paragraph()
p3.text = "\nLive App Link: http://localhost:3000\nDemo Credentials: Admin (Olumide) / Employee (Fatima)"
p3.font.name = "Segoe UI"
p3.font.size = Pt(13)
p3.font.color.rgb = TEXT_MUTED

add_footer(s10, 10, 10, is_dark=True)

# Save Presentation
prs.save("presentation.pptx")
print("Presentation presentation.pptx generated successfully!")
